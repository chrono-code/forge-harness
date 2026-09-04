#!/usr/bin/env python3
"""R1·R2·R4·R5·P1·P3 의 known-pair pptx 를 **테스트 시점에** 만든다 — 레포에 커밋하지 않는다.

🟥 R3 는 여기서 안 만든다. 선두 줄 판별이 절대좌표(LEAD_MAX_Y·LEAD_MIN_CX)와 srgbClr 런 구조에
   기대므로, 검출기 가정에 맞춰 python-pptx 로 «생성»한 픽스처는 검출기 자신을 검증 못 한다
   (자기참조). 대신 실물 구조를 그대로 2장으로 줄인 `fixture_R3_positive.pptx` ·
   `fixture_R3_negative.pptx` 를 쓴다(이 디렉터리에 실물로 있다).

슬라이드 크기는 원 코퍼스와 같은 1920pt 계열(24,384,000 × 13,716,000 EMU)로 맞춘다 — P1/P3 는
상대좌표 비교라 크기 자체는 판정에 안 들어가지만, R1/R3 의 `LEAD_MIN_CX`(12,000,000 EMU) 가
«슬라이드 절반 폭» 정도를 가정하므로 실물 규모에서 벗어나면 그 가정을 검증 못 한다.
"""
import os
from pptx import Presentation
from pptx.util import Emu
from pptx.dml.color import RGBColor

SLIDE_W, SLIDE_H = Emu(24384000), Emu(13716000)
ACCENT = RGBColor(0xFA, 0xE1, 0x00)


def _new_deck():
    prs = Presentation()
    prs.slide_width, prs.slide_height = SLIDE_W, SLIDE_H
    return prs


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _box(slide, name, x, y, w, h, text=None, accent=False):
    tb = slide.shapes.add_textbox(Emu(x), Emu(y), Emu(w), Emu(h))
    tb.name = name
    if text is not None:
        p = tb.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = text
        if accent:
            r.font.color.rgb = ACCENT
    return tb


def _notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


# ── R1 orphan-connective ─────────────────────────────────────────────────────
# 제목은 폭을 좁혀 lead 자격에서 뺀다(LEAD_MIN_CX=12,000,000 미만) — 선두 줄만 lead 로 잡히게.
def r1(path, positive):
    prs = _new_deck()
    s0 = _blank(prs)
    _box(s0, 'title0', 100000, 100000, 5000000, 500000, '첫 장')
    prev_body = ('오늘은 깃허브 사용법을 설명합니다' if not positive
                 else '오늘은 배포 절차를 설명합니다')
    _box(s0, 'body0', 100000, 2000000, 8000000, 800000, prev_body)
    s1 = _blank(prs)
    _box(s1, 'title1', 100000, 100000, 5000000, 500000, '둘째 장')
    _box(s1, 'lead1', 100000, 900000, 14000000, 900000, '그러나 깃허브는 막을 자리를 줄 뿐입니다')
    _box(s1, 'body1', 100000, 2500000, 8000000, 800000, '본문 내용입니다')
    prs.save(path)


# ── R2 enum-dropped ──────────────────────────────────────────────────────────
def r2(path, positive):
    prs = _new_deck()
    s0 = _blank(prs)
    _box(s0, 'enum0', 100000, 100000, 14000000, 900000,
         '① 속도 개선 ② 품질 향상 ③ 비용 절감')
    s1 = _blank(prs)
    _box(s1, 'title1', 100000, 100000, 14000000, 900000, '이 셋 중 무엇이 가장 중요한가')
    body = ('① 속도 ② 품질 ③ 비용을 함께 고려합니다' if not positive
            else '속도와 품질과 비용을 함께 고려합니다')
    _box(s1, 'body1', 100000, 2000000, 14000000, 900000, body)
    prs.save(path)


# ── R4 screen-heavy (자기 분포 P90, N=12) ────────────────────────────────────
# 🟥 negative 는 «값을 살짝 낮춘 아웃라이어」가 아니라 **전부 동값**으로 짠다 — N=12 처럼 작은
#    표본에서는 살짝 낮춰도 sorted 순서상 «둘째로 큰 값」이 바뀌어 p90 인덱스가 다른 정상값을
#    가리키고, 그 값이 다시 «최댓값이라 p90 초과」가 되는 자기충족 오탐이 난다(실측으로 확인).
#    전부 동값이면 어떤 값도 p90 을 **엄격히** 못 넘어(같음은 초과가 아니다) 후보가 구조적으로 0.
def r4(path, positive):
    prs = _new_deck()
    for i in range(12):
        s = _blank(prs)
        n = 500 if (positive and i == 6) else 30
        _box(s, f'body{i}', 100000, 100000, 14000000, 900000, '가' * n)
    prs.save(path)


# ── R5 read-load (신규 글자 ÷ 발화 초, 자기 분포 P90, N=12) ─────────────────
# 같은 이유로 negative 는 «전부 동일 비율」로 짠다(위 R4 와 같은 원리).
def r5(path, positive):
    """정상 슬라이드: 신규 64자 · 발화 12초(64/5.36≈12) → 비율 5.33.
    아웃라이어(positive 만, 7번째 장): 신규 150자 · 발화 5초(27/5.36≈5) → 비율 30.0."""
    prs = _new_deck()
    for i in range(12):
        s = _blank(prs)
        if positive and i == 6:
            _box(s, 'body_outlier', 100000, 100000, 14000000, 2000000,
                 '완전히새로운화면내용' * 15)  # 150자, 이전 장과 겹치지 않는 새 문자열
            _notes(s, '가' * 27)              # round(27/5.36) = 5초
        else:
            # 접두 2자리 인덱스로 매 장을 다르게 만든다(집합 차집합 비교라 «신규»로 잡히려면
            # 이전 장과 텍스트가 겹치면 안 된다) — 나머지 62자는 동일해 길이는 항상 64자다
            _box(s, f'body{i}', 100000, 100000, 14000000, 2000000, f'{i:02d}' + '가' * 62)
            _notes(s, '가' * 64)              # round(64/5.36) = 12초
    prs.save(path)


# ── P1 build-jitter (같은 이름 도형, 연속 4장) ───────────────────────────────
def p1(path):
    """4장 · 이름 «gate_bar» 공유. 쌍(1→2)=임계 안 소폭 이동(뜸) · (2→3)=임계 밖 큰 이동(안 뜸) ·
    (3→4)=무이동(안 뜸). compared=3 이어야 «비교 안 함»이 아니다."""
    prs = _new_deck()
    xs = [1000000, 1030000, 6030000, 6030000]  # +30,000(JIT 200,000 안) · +5,000,000(밖) · +0
    for x in xs:
        s = _blank(prs)
        _box(s, 'gate_bar', x, 2000000, 1000000, 500000, '게이트 막대')
    prs.save(path)


# ── P3 adjacency (오른끝 ↔ 왼끝) ─────────────────────────────────────────────
def p3(path):
    """slide1 = 임계 안 어긋남(뜸) · slide2 = 정확히 맞닿음(0, 안 뜸) + 임계 밖 큰 틈(안 뜸)."""
    prs = _new_deck()
    s1 = _blank(prs)
    _box(s1, 'stub_a', 1000000, 2000000, 1000000, 500000, 'A')     # 오른끝 2,000,000
    _box(s1, 'bar_b', 2020000, 2000000, 1000000, 500000, 'B')      # 왼끝 2,020,000 · 틈 20,000(<ALN)
    s2 = _blank(prs)
    _box(s2, 'stub_c', 1000000, 2000000, 1000000, 500000, 'C')     # 오른끝 2,000,000
    _box(s2, 'bar_d', 2000000, 2000000, 1000000, 500000, 'D')      # 왼끝 2,000,000 · 틈 0 (맞닿음)
    _box(s2, 'stub_e', 1000000, 4000000, 1000000, 500000, 'E')     # 오른끝 2,000,000
    _box(s2, 'bar_f', 3000000, 4000000, 1000000, 500000, 'F')      # 왼끝 3,000,000 · 틈 1,000,000(>ALN)
    prs.save(path)


# ── P1/P3 --baseline 델타 쌍 ──────────────────────────────────────────────────
def geometry_baseline_pair(base_path, edited_path):
    """base = 깨끗한 2장(후보 0). edited = 2장째 도형 하나를 소폭(30,000 EMU) 옮긴 사본
    (새 P1 후보 정확히 1건)."""
    prs = _new_deck()
    s1 = _blank(prs)
    _box(s1, 'gate_bar', 1000000, 2000000, 1000000, 500000, '게이트')
    s2 = _blank(prs)
    _box(s2, 'gate_bar', 1000000, 2000000, 1000000, 500000, '게이트')
    prs.save(base_path)

    prs2 = _new_deck()
    t1 = _blank(prs2)
    _box(t1, 'gate_bar', 1000000, 2000000, 1000000, 500000, '게이트')
    t2 = _blank(prs2)
    _box(t2, 'gate_bar', 1030000, 2000000, 1000000, 500000, '게이트')  # +30,000 EMU
    prs2.save(edited_path)


def build_all(outdir):
    os.makedirs(outdir, exist_ok=True)
    r1(os.path.join(outdir, 'r1_pos.pptx'), True)
    r1(os.path.join(outdir, 'r1_neg.pptx'), False)
    r2(os.path.join(outdir, 'r2_pos.pptx'), True)
    r2(os.path.join(outdir, 'r2_neg.pptx'), False)
    r4(os.path.join(outdir, 'r4_pos.pptx'), True)
    r4(os.path.join(outdir, 'r4_neg.pptx'), False)
    r5(os.path.join(outdir, 'r5_pos.pptx'), True)
    r5(os.path.join(outdir, 'r5_neg.pptx'), False)
    p1(os.path.join(outdir, 'p1.pptx'))
    p3(os.path.join(outdir, 'p3.pptx'))
    geometry_baseline_pair(os.path.join(outdir, 'geo_base.pptx'),
                            os.path.join(outdir, 'geo_edited.pptx'))
    return outdir


if __name__ == '__main__':
    import sys, tempfile
    d = sys.argv[1] if len(sys.argv) > 1 else tempfile.mkdtemp(prefix='preprep_fixtures_')
    build_all(d)
    print('built into', d)
