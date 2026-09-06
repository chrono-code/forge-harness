#!/usr/bin/env python3
"""L13 — «N 페이지» 리터럴 장 번호 참조가 실재하는 장을 가리키나 (2026-09-06 신설).

WHY: 장을 하나 빼면 그 뒤 번호가 통째로 당겨진다. preprep 의 이웃 레인들은 **말로 가리키는
것**만 본다 — L5/L8 은 화면 재사용·시각 지시를, R1~R5 는 접속어·소유권을. **숫자로 가리키는
것은 어느 레인도 안 봤다.** 실사고 계기: if(kakao) 덱에서 42장을 빼자 43~121 이 한 칸씩
당겨졌고, 「깨진 참조가 어디냐」를 손으로 훑어야 했다(2026-09-06, 운영자 지적).

🟥 이 레인의 난점은 검출이 아니라 **분류**다. 같은 «숫자+단위»가 세 가지 다른 것을 가리킨다 —
실측 코퍼스(그 덱 121장)에서 셋이 **전부** 나왔다:

    SLIDE     「40p 부제가 …」          → 진짜 장 참조. 이것만 검사 대상이다
    SECTION   「3장에서 게이트부터」      → 슬라이드가 아니라 **섹션**(그 덱의 sec-guide 「3. 어떻게 만들었나」)
    CITATION  「arXiv:2606.05647 (4p)」 → 인용 **논문 쪽수**

셋을 안 가르면 오탐 덩어리가 된다. 그래서 판별자는 코퍼스에서 기계로 뽑는다:
  · SECTION  — 덱의 섹션 라벨(「N. 제목」)에서 N 집합을 모으고, 「N장」 형태이면서 그 집합에
              속하면 섹션이다. 라벨이 없으면 이 분기는 **비활성**이고 그 사실을 노트에 적는다.
  · CITATION — 같은 줄에 인용 표지(arXiv/doi/『』/pp./4자리 연도)가 있으면 인용 쪽수다.

🟥 **판정은 «존재»만 한다** — «그 번호가 맞는 장을 가리키나»는 의미 판단이라 안 본다
(§Mechanization Boundary: 기록의 속성만 단언한다). 그래서 finding 은 **범위 밖 참조**
하나뿐이고, 자기 자신 참조·정상 참조는 노트로만 나간다.

🟥 **참조 0건은 «통과»가 아니다** — 추출이 죽어도 0이 나온다. 0이면 UNMEASURED 로 적는다.
"""
import os
import re

# 「40p」 「40 페이지」 「40쪽」 「3장에서」…
REF = re.compile(r'(?<![0-9A-Za-z])(\d{1,3})\s*(p\b|페이지|쪽(?!수)|장(?=에서|의|을|부터|에))')
# 🟥 `\s+` 가 하중을 진다 — 초판은 `\s*` 였고 실측 코퍼스의 «62.5%» 를 «62번 섹션»으로
#    먹었다. 그러면 진짜 「62장」 참조가 조용히 SECTION 으로 분류돼 사라진다.
#    섹션 라벨은 점 뒤에 공백이 있다(「3. 어떻게 만들었나」). 소수점은 없다.
SECTION_LABEL = re.compile(r'^\s*(\d{1,2})\.\s+\S')
SECTION_MAX = 20      # 섹션이 20을 넘는 덱은 없다 — 넘으면 라벨이 아니다
SECTION_MIN_SLIDES = 2  # 섹션 라벨은 여러 장에 반복된다(한 장짜리는 라벨이 아니다)
CITE = re.compile(r'arXiv|doi|『|』|pp\.|\b(19|20)\d{2}\b', re.I)


def _texts(slide):
    scr, note = [], ''
    for sh in slide.shapes:
        if sh.has_text_frame:
            t = ' '.join(r.text for p in sh.text_frame.paragraphs for r in p.runs)
            if t.strip():
                scr.append(t)
    if slide.has_notes_slide:
        note = slide.notes_slide.notes_text_frame.text or ''
    return scr, note


def section_numbers(prs):
    """섹션 라벨(「N. 제목」)에서 N 을 모은다. 못 모으면 빈 집합 — 그 사실이 노트에 적힌다.

    세 조건을 **함께** 건다: 점 뒤 공백(소수점 배제) · N ≤ SECTION_MAX · 여러 장에 반복.
    하나만 걸면 실측에서 뚫린다(위 주석의 62.5% 사례).
    """
    seen = {}
    for i, s in enumerate(prs.slides, 1):
        for sh in s.shapes:
            if not sh.has_text_frame:
                continue
            for p in sh.text_frame.paragraphs:
                line = ''.join(r.text for r in p.runs).strip()
                m = SECTION_LABEL.match(line)
                if m and len(line) < 40:
                    n = int(m.group(1))
                    if n <= SECTION_MAX:
                        seen.setdefault(n, set()).add(i)
    return {n for n, slides in seen.items() if len(slides) >= SECTION_MIN_SLIDES}


def classify(num, unit, line, secs):
    if CITE.search(line):
        return 'CITATION'
    if unit.startswith('장') and num in secs:
        return 'SECTION'
    return 'SLIDE'


def analyze(prs):
    total = len(prs.slides)
    secs = section_numbers(prs)
    rows = []
    for i, s in enumerate(prs.slides, 1):
        scr, note = _texts(s)
        for where, chunks in (('화면', scr), ('대본', note.split('\n'))):
            for line in chunks:
                for m in REF.finditer(line):
                    n, unit = int(m.group(1)), m.group(2)
                    kind = classify(n, unit, line, secs)
                    rows.append((i, where, n, unit, kind, line.strip()[:70]))
    return total, secs, rows


def scan(cfg, root):
    deck_s = (cfg.get('surfaces_by_id') or {}).get('built_deck')
    if not deck_s:
        return [], ['L13 slide-refs : built_deck 미선언 — NOT_CONFIGURED (0 아님)']
    path = os.path.normpath(os.path.join(root, os.path.expanduser(deck_s['path'])))
    if not os.path.exists(path):
        return [], [f'L13 slide-refs : built_deck 실물 없음({path}) — UNMEASURED (0 아님)']
    try:
        from pptx import Presentation
        prs = Presentation(path)
        total, secs, rows = analyze(prs)
    except Exception as e:
        return [], [f'L13 slide-refs : 계기 오류({type(e).__name__}: {e}) — UNMEASURED (0 아님)']

    per = {}
    for _, _, _, _, k, _ in rows:
        per[k] = per.get(k, 0) + 1
    notes = [f'L13 slide-refs : {total}장 · 참조 {len(rows)}건 '
             f'(SLIDE {per.get("SLIDE",0)} · SECTION {per.get("SECTION",0)} · CITATION {per.get("CITATION",0)}) '
             f'· 섹션 번호 집합 {sorted(secs) if secs else "없음 — SECTION 분기 비활성"}']
    if not rows:
        notes.append('   ⚠️ 참조 0건 — 「깨끗하다」가 아니라 **UNMEASURED** 다(추출이 죽어도 0이 나온다). '
                     '덱에 번호 참조가 원래 없으면 정상이고, 그 사실을 여기서 확인해라')
        return [], notes

    findings = []
    for i, where, n, unit, kind, line in rows:
        if kind != 'SLIDE':
            continue
        if n < 1 or n > total:
            findings.append((f'{path}#s{i}', 'use',
                             f'L13 장 번호 참조가 범위 밖이다 — s{i} [{where}] →{n} (총 {total}장) : {line}'))
            notes.append(f'   ❌ s{i:3} [{where}] →{n:3}  범위 밖(총 {total}) : {line}')
        elif n == i:
            notes.append(f'   ⚠️ s{i:3} [{where}] →{n:3}  자기 자신 참조 : {line}')
        else:
            notes.append(f'   ·  s{i:3} [{where}] →{n:3}  : {line}')
    return findings, notes
