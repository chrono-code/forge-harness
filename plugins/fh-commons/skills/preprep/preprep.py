#!/usr/bin/env python3
"""발표 준비 하네스 v0.1 — 발표 자산군 전체(덱 + 기술문서 + Q&A + 초록)를 한 번에 얼라인한다.

챔버 런 #14 (`tracks/_chamber/multisurface-reading-harness/`) 의 EMIT 산출물.
「다표면 읽기」는 이 하네스의 **한 렌즈(L2)** 이지 전체 이름이 아니다(운영자 위계 정정 2026-08-20).

레인:
  L1 canon      — 원장 「🚫 인용금지 / ⚠️ 조건필수」를 **전 표면**에 건다 (한글 낭독형 포함)
  L2 unit-ref   — 표면 쌍의 단위 참조 정합 (대상의 «은퇴 선언 어휘»를 존중 — 안 하면 20% FP)
  L3 inventory  — 선언된 표면이 실물로 있나 (부재 = UNMEASURED, **0 이 아니다**)
  L4 getput     — 왕복 no-op 컨트롤. **v0.1 미배선 → NOT_WIRED 로 말한다**(거짓 초록 금지)
  L12 diagram   — 타입 JSON 도해가 «지금 JSON» 에서 validate 를 거쳐 구워졌나 (lane_diagram.py · 2026-09-05)

exit 0 = 전 레인 통과 · 1 = 발견 있음 · 2 = 계기 오류/미측정으로 판정 불가 (PASS 아님)

usage: preprep.py [surfaces.yaml] [--lane L1,L2,...]
"""
import sys, os, re, ast, zipfile, yaml

HERE = os.path.dirname(os.path.abspath(__file__))
ANSI_OK, ANSI_BAD, ANSI_UNK = '🟢', '🟥', '⚠️'

def load(path):
    with open(path, encoding='utf-8') as f: return yaml.safe_load(f)

def resolve(root, p):
    return os.path.normpath(os.path.join(root, os.path.expanduser(p)))

# 🟥 2026-09-04 수리(신호 §3-2) — canon_terms/jargon_terms 는 이제까지 **HERE(이 스크립트
#    위치)** 기준으로만 풀렸다. 자산군 쪽(surfaces.yaml 이 있는 root)에 실제 목록 파일을 두면
#    못 찾았고, 실사용에서 절대경로를 넣어서야 돌았다(문서에도 없던 함정).
#    ⇒ root 를 먼저 보고, 없을 때만 HERE(스킬 배포본에 딸린 example 계열)로 물러난다 —
#    스킬을 설치 위치와 무관하게 그대로 돌릴 수 있어야 하므로 폴백 자체는 남긴다. 폴백을
#    썼다는 사실은 출력에 «(skill-dir fallback)» 로 남긴다(조용히 넘기지 않는다).
def resolve_declared(root, value):
    """(경로, skill-dir 폴백 여부). 절대경로면 그대로. root 에 없으면 HERE 로 시도한다."""
    if os.path.isabs(value):
        return value, False
    p_root = resolve(root, value)
    if os.path.exists(p_root):
        return p_root, False
    p_here = resolve(HERE, value)
    if os.path.exists(p_here):
        return p_here, True
    return p_root, False   # 둘 다 없다 — root 기준 경로를 대서 에러 메시지가 뭘 찾았는지 보이게 한다

# 🟥 2026-08-25 신설 — 선언 경로가 판번호로 하드코딩돼 stale 해지는 클래스(G2).
#    실측: surfaces.yaml 이 v4.6 을 가리키는 동안 실제 산출물은 v7.5 였다(12판 차이).
#    ⚠️ **최신본으로 자동 추종하지 않는다** — 옛 판을 «일부러» 재는 경우를 뺏기 때문이다.
#    대신 ⓐ 무엇을 쟀는지 항상 출력하고 ⓑ 선언 ≠ 디스크 최신본이면 그 자체를 발견으로 낸다.
def version_drift(path):
    import glob
    d, b = os.path.dirname(path), os.path.basename(path)
    m = re.match(r'(.*?_v)(\d+)\.(\d+)(\.pptx)$', b)
    if not m: return None
    cand = []
    for f in glob.glob(os.path.join(d, m.group(1) + '*' + m.group(4))):
        mm = re.match(r'.*?_v(\d+)\.(\d+)\.pptx$', os.path.basename(f))
        if mm: cand.append(((int(mm.group(1)), int(mm.group(2))), os.path.basename(f)))
    if not cand: return None
    newest = max(cand)
    here = (int(m.group(2)), int(m.group(3)))
    if newest[0] > here:
        return (b, newest[1])
    return None

# ── 표면 읽기 어댑터 ───────────────────────────────────────────────────────────
def read_text(path, kind):
    """표면의 «읽을 수 있는 텍스트». 못 읽으면 예외 — 호출자가 UNMEASURED 로 계상한다.
    🟥 pptx 는 XML 텍스트 레이어만 읽는다. **이미지에 구워진 텍스트는 구조적으로 안 보인다**
       (실측 2026-08-18) — L1 이 그 한계를 출력에 명시한다. 미측정을 0 으로 접지 않는다."""
    if kind in ('markdown', 'html', 'note_map'):
        return open(path, encoding='utf-8', errors='replace').read()
    if kind == 'pptx':
        out = []
        with zipfile.ZipFile(path) as z:
            for n in z.namelist():
                if re.match(r'ppt/(slides|notesSlides)/\w+\.xml$', n):
                    out.append(re.sub(r'<[^>]+>', ' ', z.read(n).decode('utf-8', 'replace')))
        return '\n'.join(out)
    if kind == 'figure_source':
        # 🟥 2026-08-21 신설 — 이 파일 헤더가 「이미지에 구워진 텍스트는 구조적으로 안 보인다」고
        #   선언해 뒀는데, 그 사각에서 같은 날 **두 건이 물었다**: 도해 안의 「초록이 는 게
        #   아닙니다」(오타 + 덱이 버린 은어)와 도해 전면의 「계열」(덱은 「모델/종류」로 갔다).
        #   둘 다 화면에는 떠 있고 텍스트 레이어엔 없어서, 눈으로 보기 전엔 아무 레인도 못 잡았다.
        #   ⇒ **그림은 못 읽지만 그림을 그리는 소스는 읽을 수 있다.** 생성기의 문자열 리터럴만
        #     뽑아 표면으로 세운다. 이러면 L1(원장 금칙)·L5(조어)가 도해까지 덮는다.
        #   ⚠️ 이건 «그림을 읽는 것»이 아니다 — 생성기를 안 거친 그림(손으로 넣은 png)은
        #     여전히 안 보인다. 사각이 좁아진 것이지 닫힌 게 아니다.
        src = open(path, encoding='utf-8', errors='replace').read()
        lits = re.findall(r'"([^"\n]{2,})"|\'([^\'\n]{2,})\'', src)
        return '\n'.join(a or b for a, b in lits)
    if kind == 'diagram_source':
        # 🟥 2026-09-05 신설 — 타입 있는 도해 JSON(archify 스키마). figure_source 와 같은 이유로
        #   «그림은 못 읽지만 그림을 그리는 소스는 읽는다». 노드 label/sublabel/tag · 간선 label ·
        #   레인/구간/그룹 label · 카드 문장만 뽑는다(id·type 같은 기계 낱말은 L1/L5 에 안 태운다).
        #   굽기 정합(지금 JSON 의 그림인가 · validate · 해상도)은 L12(lane_diagram.py) 몫이다.
        import json as _json
        d = _json.load(open(path, encoding='utf-8'))
        out = []
        for key in ('lanes', 'phases', 'groups', 'nodes', 'edges', 'participants', 'messages', 'states', 'transitions', 'components', 'relationships', 'stages', 'flows'):
            for it in d.get(key) or []:
                if isinstance(it, dict):
                    for f in ('label', 'sublabel', 'tag', 'note'):
                        if it.get(f): out.append(str(it[f]))
        for c in d.get('cards') or []:
            if isinstance(c, dict):
                if c.get('title'): out.append(str(c['title']))
                out += [str(x) for x in (c.get('items') or [])]
        m = d.get('meta') or {}
        for f in ('title', 'subtitle'):
            if m.get(f): out.append(str(m[f]))
        for v in m.get('views') or []:
            if isinstance(v, dict):
                for f in ('label', 'note'):
                    if v.get(f): out.append(str(v[f]))
        return '\n'.join(out)
    if kind == 'keynote':
        raise RuntimeError('keynote: 사람만 읽는 면 — 기계 어댑터 없음')
    raise RuntimeError(f'unknown kind: {kind}')

def note_map_ids(path):
    src = open(path, encoding='utf-8').read()
    m = re.search(r'^NOTE_MAP\s*=\s*\{', src, re.M)
    if not m: raise RuntimeError('NOTE_MAP 미발견')
    start, depth = m.end() - 1, 0
    for i in range(start, len(src)):
        if src[i] == '{': depth += 1
        elif src[i] == '}':
            depth -= 1
            if depth == 0:
                return {v[0] for v in ast.literal_eval(src[start:i+1]).values() if v[0]}
    raise RuntimeError('NOTE_MAP 괄호 불균형')

def manuscript_ids(path, pattern, retire_markers):
    """🟥 은퇴 선언을 **헤딩에서만** 찾던 초판은 4장을 살아있는 것으로 계상했다.
    실측(2026-08-20): S19 만 헤딩에 선언이 있고, S20-c·S20-d 는 **인용줄**에,
    S8·S18-b 는 **아예 선언이 없다**. ⇒ 판정면을 «장 블록»으로 넓힌다.
    선언이 없는 장은 그대로 남아 진짜 발견이 된다(그게 이 레인이 잡아야 할 것이다).
    (블라인드 challenger 가 소스 직독으로 지목 → 거버너가 grep 으로 확정)"""
    src = open(path, encoding='utf-8').read()
    ids = []
    heads = list(re.finditer(pattern + r'([^\n]*)', src, re.M))
    for i, m in enumerate(heads):
        end = heads[i+1].start() if i+1 < len(heads) else len(src)
        block = src[m.start():end][:1500]          # 선언은 장 앞부분에 온다
        if any(t in block for t in retire_markers): continue
        ids.append(m.group(1))
    if not ids: raise RuntimeError('섹션 헤딩 0건 — 포맷 불일치')
    return set(ids)

# ── use vs mention ─────────────────────────────────────────────────────────────
# 🟥 이 구분이 L1 의 하중선이다. 첫 실행에서 40건이 나왔는데 손검증 3/3 이 **«쓰지 마라»고 적은
#    줄**이었다 — 원장 자신이 기록한 자기참조 결함(«세는 대상 안에 출처 문서가 들어 있어, 계기가
#    자기 정의를 착지로 셌다»)이 이 계기에 그대로 재현됐다.
# 🟥 그렇다고 **걸러서 지우지 않는다** — 죽은 필터가 진짜를 죽인다. 두 부류로 **갈라서 낸다**.
_HIST = re.compile(r'^\s*>\s*(🆕|🟥|⚠️|✅|🔵|🎯|📌|\*\*20\d\d-|운영자[ :]|원인\(|기전|\*"|정본 ?=)')

MENTION_CUES = ('🚫', '금지', '대신 쓸 것', '폐기', 'stale', '반증', '원장', '오기',
                '쓰지 않는다', '안 쓴다', '등재', '철회', '죽었', '틀린',
                # 손검증에서 나온 추가 큐 (2026-08-20, USE 표본 손대조) — 정정 서술의 어휘
                '교체', '원문', '건너뛴', '구값', '대체', '재측정 필수', '→')

def classify(line):
    """USE(살아있는 사용) / MENTION(그것을 말하지 말라고 적은 줄). 보수적으로 — 애매하면 USE.

    🟥 초판은 `>` 로 시작하면 무조건 MENTION 이었다. **이 코퍼스는 낭독도 `>` 다**(1232줄).
    그래서 원고 표면의 L1 판정 건수가 **0** 이었다 — 가장 중요한 표면을 통째로 못 봤고,
    출력은 초록으로 보였다. 「죽은 필터가 진짜를 죽인다」를 이 파일이 경고해놓고 저질렀다.
    (블라인드 페르소나 2인이 각각 지목 · 거버너 자력 적발 0 · 실행 한 번으로 확정)
    ⇒ 판별자는 인용부호가 아니라 **주석 어휘**다. `>` 는 이제 아무 정보도 아니다."""
    t = line.strip().lstrip('> ').strip()
    if t.startswith('#'): return 'MENTION'          # 소스 주석 — 청중이 안 본다
    # 대본 안의 «*(화면: …)*» 는 낭독이 아니라 화면 지시다. 원장이 그 둘을 다르게 규정한다
    # (절대수: 화면 허용 · 낭독 금지) — 표면 단위로만 갈면 이 줄이 오탐으로 남는다(실측).
    if t.startswith('*(화면') or t.startswith('(화면'): return 'MENTION'
    if t.startswith('>'): return 'MENTION'          # 인용/주석 블록
    if t.startswith('|') and any(c in t for c in MENTION_CUES): return 'MENTION'   # 원장 표 행
    if any(c in t for c in MENTION_CUES): return 'MENTION'
    return 'USE'

def hits(text, lit, speech_marker=None):  # yields (cls, line_no, line)
    """speech_marker 가 선언되면 **그 표시가 여는 블록 안**만 USE 후보다 — 대상 자신의 어휘로
    범위를 좁힌다(은퇴 표기와 같은 원리, 그쪽은 실측으로 20% FP 를 없앴다)."""
    in_speech = speech_marker is None
    for no, ln in enumerate(text.split('\n'), 1):
        if speech_marker:
            # 🟥 «줄 안에 있으면» 이 아니라 «줄이 그것으로 시작하면». 2026-08-20 실측:
            # 이 코퍼스의 앞머리에 `각 장의 \`🗣\` = 발표자가 말할 것` 이라는 **범례**가 있어,
            # 기호를 설명하는 줄이 기호로 읽히고 그 뒤 앞머리 75줄이 통째로 «낭독»이 됐다.
            # 실제 낭독 블록은 항상 줄머리에 표시자가 온다(L98 등) — 그게 판별자다.
            if ln.lstrip().startswith(speech_marker): in_speech = True
            elif ln.startswith('###') or ln.startswith('🖥') or ln.startswith('⏱'): in_speech = False
            # 🟥 이력 블록이 열리면 낭독 블록은 끝난 것이다. 이 코퍼스는 개정 이력을 낭독과
            # **같은 `>` 인용 블록 안**에 이어 붙이고(구판 문장을 통째로 인용하는 일이 흔하다),
            # 그래서 「그럼 이 적대검증이 … 보여 드리겠습니다」 같은 **폐기된 낭독**이 살아있는
            # 낭독으로 계상됐다. 표시자로 열고 이력 표시로 닫는다 — 이 판별은 이 파일이 이미
            # L2 에서 쓰는 «대상 자신의 어휘로 범위를 좁힌다»와 같은 원리다. (2026-08-20 실측)
            elif _HIST.match(ln): in_speech = False
        if lit in ln:
            cls = classify(ln)
            if cls == 'USE' and not in_speech: cls = 'mention'
            yield cls, no, ln.strip()

# ── L1 canon ──────────────────────────────────────────────────────────────────
def lane_canon(cfg, root, texts, surf_meta):
    findings, notes = [], []
    ct_cfg = cfg.get('canon_terms')
    if not ct_cfg:
        return [], ['L1 canon : canon_terms 미선언 — NOT_CONFIGURED (0 아님)']
    ct_path, ct_fb = resolve_declared(root, ct_cfg)
    if ct_fb:
        notes.append(f'L1 canon : canon_terms 를 root 에서 못 찾아 스킬 위치로 폴백 '
                     f'(skill-dir fallback) — {ct_path}')
    # 🟥 2026-09-04 수리(신호 §3-3) — 파일 부재가 여기서 `FileNotFoundError` 트레이스백으로
    #    죽었다. 이 스킬 자신의 교리가 「부재 ↔ 0 을 가른다 · UNMEASURED 로 낸다」인데 자기
    #    코드가 그걸 어기고 있었다([[feedback_rule_misdescribes_its_own_machine]] 과 같은 족).
    #    ⇒ UNMEASURED 로 낮춰 잡고 **나머지 레인은 계속 돈다**(L1 하나가 죽는다고 전체가
    #    안 죽는다 — main() 은 이 함수를 감싸지 않으므로 여기서 반드시 잡아야 한다).
    try:
        terms = load(ct_path)
    except Exception as e:
        return [], [f'L1 canon : canon_terms 부재({ct_path}) — UNMEASURED (0 아님): '
                    f'{type(e).__name__}: {e}']

    # L1-a 원장 ↔ 기계 목록 커버리지 (이 하네스가 잡는 결함이 자기 자신에게도 난다)
    ledger = resolve(root, cfg['canon_ledger'])
    try:
        rows = [l for l in open(ledger, encoding='utf-8') if l.startswith('| **') ]
        n_terms = (len(terms.get('banned', [])) + len(terms.get('conditional', []))
                   + len(terms.get('retired', [])))
        notes.append(f'L1-a 원장 표 행 {len(rows)} ↔ 기계 항목 {n_terms} — '
                     + ('커버 비율은 판정 안 한다(행≠항목 1:1 아님). 갈림 감시용 수치' ))
    except Exception as e:
        notes.append(f'L1-a UNCOVERED: 원장을 못 읽었다 ({e})')

    for sid, (text, kind) in texts.items():
        for entry in terms.get('banned', []):
            if entry.get('spoken_only') and not surf_meta.get(sid, {}).get('spoken'):
                continue                              # 낭독 표면에서만 판정하는 항목
            for lit in entry['literals'] + entry.get('ko_spoken', []):
                if not lit: continue
                for cls, no, ln in hits(text, lit, surf_meta.get(sid, {}).get('marker')):
                    # 리터럴이 너무 일반적일 때: 같은 줄에 맥락 어휘가 있어야 판정한다
                    #  (실측: «0.44» 가 슬라이드 y좌표 0.44 를 세 번 잡았다)
                    if cls == 'USE' and entry.get('context_required') and \
                       not any(c in ln for c in entry['context_required']):
                        cls = 'mention'
                    findings.append((f'{sid}:{no}', 'BANNED' if cls == 'USE' else 'mention',
                                     entry['id'], lit, (entry['why'] + ' ▸ ' + ln[:110]) if cls == 'USE' else ln[:110]))
        # ── L1-c 폐어(retired) — «리뷰가 걷어낸 낱말이 되돌아왔나» ────────────────
        # 🟥 실사고: 리뷰가 명시적으로 걷어내라 한 「재다」가 **나흘 뒤** 재유입됐고
        #    **아무 검사도 안 울렸다**(운영자가 잡았다). `jargon_terms` 는 조어→풀이를,
        #    `banned` 는 수치 인용 금지를 다뤄서 «걷어낸 말»을 담을 자리가 없었다.
        # 🟥 왜 새 레인이 아니라 여기인가: 기계는 «리터럴을 전 표면에서 찾기»로 같은데
        #    **처방이 다르다**(폐어=대체어를 써라 / 금지인용=그 재측정을 대라). 태그와 처방만
        #    가르면 «판정은 맞고 처방이 틀린 게이트»를 피하면서 레인 하나를 안 늘린다.
        # 🟥 그리고 USE↔MENTION 이 여기서 하중을 진다 — 폐어를 «걷어냈다»고 적은 원장 자신이
        #    그 낱말을 담는다. 못 가르면 원장이 자기를 위반으로 신고한다.
        for entry in terms.get('retired', []):
            for lit in entry['literals'] + entry.get('ko_spoken', []):
                if not lit: continue
                for cls, no, ln in hits(text, lit, surf_meta.get(sid, {}).get('marker')):
                    if cls != 'USE':
                        findings.append((f'{sid}:{no}', 'mention', entry['id'], lit, ln[:110]))
                        continue
                    prov = ' · '.join(x for x in (entry.get('retired_by'), entry.get('retired_at')) if x)
                    rep = entry.get('replacement')
                    msg = (entry.get('why', '리뷰가 걷어낸 낱말')
                           + (f" ▸ 대신: «{rep}»" if rep else " ▸ 🟥 대체어 미기재 — 처방 없는 판정이다")
                           + (f" ({prov})" if prov else '') + ' ▸ ' + ln[:110])
                    findings.append((f'{sid}:{no}', 'RETIRED', entry['id'], lit, msg))
        for entry in terms.get('conditional', []):
            for lit in entry['literals']:
                if lit and lit in text:
                    for cls, no, ln in hits(text, lit, surf_meta.get(sid, {}).get('marker')):
                        if cls != 'USE': continue
                        # 조건 어휘는 **같은 줄**에서 본다 — 문서 전체에서 보면 언제나 통과한다
                        if not any(r in ln for r in entry['require_nearby']):
                            findings.append((f'{sid}:{no}', 'UNCONDITIONED', entry['id'], lit,
                                             entry['why'] + ' ▸ ' + ln[:110]))
    return findings, notes

# ── L2 unit-ref ───────────────────────────────────────────────────────────────
def lane_unitref(cfg, root, surf):
    findings, notes = [], []
    for a, b in cfg.get('pairs', []):
        sa, sb = surf[a], surf[b]
        try:
            ids_a = manuscript_ids(resolve(root, sa['path']), sa['unit_pattern'],
                                   sa.get('retire_markers', []))
            ids_b = note_map_ids(resolve(root, sb['path']))
        except Exception as e:
            notes.append(f'L2 INSTRUMENT ERROR ({a}↔{b}): {e} — 미측정이지 0 이 아니다')
            continue
        only_a, only_b = sorted(ids_a - ids_b), sorted(ids_b - ids_a)
        if only_a: findings.append((a, 'DEAD-REF', b, ', '.join(only_a),
                                    f'{a} 만 참조 — 독자에겐 살아 보이지만 {b} 에 없다'))
        if only_b: findings.append((b, 'ORPHAN', a, ', '.join(only_b),
                                    f'{b} 에만 있다 — 원고 없는 단위가 빌드된다'))
        if not only_a and not only_b:
            notes.append(f'L2 {a}↔{b}: 드리프트 0 ({len(ids_a)} 단위 일치)')
    return findings, notes

# ── L4 getput (왕복 no-op 컨트롤) ───────────────────────────────────────────────
# BX/lens 의 **GetPut 법칙**에서 빌려왔다: 뷰를 안 고쳤으면 소스도 안 바뀌어야 한다.
# 🟥 지금까지 있던 것은 PutGet 하나뿐이었다(사람이 고친 뒤 «반영됐나» 28/28). GetPut 이 없으면
#    **왕복기가 항상 무언가 바꾸는 상태여도 초록이 난다** — 컨트롤 없는 계기다.
# 기계만으로 도는 축을 쓴다: 원고 🗣 → [빌드] → pptx 노트 → 정규화 → 원고와 대조.
#    사람 편집이 0 인 구간이므로 **차이는 전부 왕복기 손실**이다.
# ⚠️ Keynote 경로(osascript)는 사람 면이라 여기서 안 쓴다 — 그쪽은 영구 UNMEASURED.
def lane_getput(cfg, root, surf):
    import importlib.util
    findings, notes = [], []
    scripts_dir = os.path.normpath(os.path.join(root, '..', 'scripts'))
    pull = os.path.join(scripts_dir, 'ifkakao_pull_notes.py')
    if not os.path.exists(pull):
        return [], ['L4 NOT_WIRED — 왕복기 부재 (UNMEASURED, 통과 아님)']
    try:
        spec = importlib.util.spec_from_file_location('pn', pull)
        pn = importlib.util.module_from_spec(spec); spec.loader.exec_module(pn)
        tags = pn.builder_tags(); nm = pn.note_map()
        draft, _sc = pn.draft_speech()
    except Exception as e:
        return [], [f'L4 INSTRUMENT ERROR: {e} — 미측정이지 0 이 아니다']

    deck = resolve(root, surf['built_deck']['path'])
    notes.append(f'L4 잰 산출물: {os.path.basename(deck)}')   # 무엇을 쟀는지가 판정의 일부다
    _vd = version_drift(deck)
    if _vd:
        findings.append(('roundtrip', 'DECLARED-STALE', 'surfaces.yaml',
                         f'선언 {_vd[0]} 인데 디스크 최신본은 {_vd[1]}',
                         '의도적으로 옛 판을 재는 것이면 사유를 선언에 적어라. '
                         '아니면 surfaces.yaml 의 built_deck 을 갱신해라 — 이 레인은 선언을 따른다'))
    try:
        # 🟥 2026-08-25 수리 — notesSlide 「파일 번호」를 「발표 순번」으로 쓰면 안 된다.
        #    OOXML 은 그 둘을 관계(rels)로 잇지 파일명으로 잇지 않는다. 실측: 이 덱의 노트는
        #    3~55 번을 쓰는데 슬라이드는 1~53 이고, slide19→notesSlide18 처럼 어긋난다.
        #    옛 코드는 notes_xml[파일번호] 를 notes_xml[순번] 으로 조회해 **전건을 밀린 채** 쟀다.
        #    매핑은 sldIdLst(순서) → presentation.rels(슬라이드 파일) → slides/_rels(노트 파일).
        def _para_text(raw):
            paras = []
            for pm in re.finditer(r'<a:p>(.*?)</a:p>', raw, re.S):
                t = ''.join(re.findall(r'<a:t>([^<]*)</a:t>', pm.group(1)))
                if t.strip(): paras.append(t)
            return '\n'.join(paras)

        notes_xml = {}
        no_note = []
        with zipfile.ZipFile(deck) as z:
            pres  = z.read('ppt/presentation.xml').decode('utf-8', 'replace')
            prels = z.read('ppt/_rels/presentation.xml.rels').decode('utf-8', 'replace')
            rid2t = dict(re.findall(r'Id="(rId\d+)"[^>]*Target="([^"]+)"', prels))
            order = re.findall(r'<p:sldId [^>]*r:id="(rId\d+)"', pres)
            for pos, rid in enumerate(order, start=1):
                tgt = rid2t.get(rid)
                if not tgt:
                    no_note.append(pos); continue
                base = os.path.basename(tgt)
                srel = 'ppt/slides/_rels/%s.rels' % base
                if srel not in z.namelist():
                    no_note.append(pos); continue
                nm_ = re.search(r'Target="([^"]*notesSlide[^"]*)"',
                                z.read(srel).decode('utf-8', 'replace'))
                if not nm_:
                    no_note.append(pos); continue
                npath = 'ppt/notesSlides/' + os.path.basename(nm_.group(1))
                if npath not in z.namelist():
                    no_note.append(pos); continue
                notes_xml[pos] = _para_text(z.read(npath).decode('utf-8', 'replace'))
    except Exception as e:
        return [], [f'L4 INSTRUMENT ERROR: pptx 노트 추출 실패 — {e} (UNMEASURED)']

    # 🟥 컨트롤이 살아 있나부터 본다 — 노트가 0장이면 «차이 0» 은 통과가 아니라 미측정이다
    if not notes_xml:
        return [], ['L4 UNMEASURED — pptx 에 노트 슬라이드 0장. «차이 0» 아님']

    # 판정에서 빼는 단위는 **선언으로** 받는다 — 코드에 박으면 그게 다음 드리프트다
    skip_units = set(cfg.get('getput_skip_units', []))
    compared = skipped = unmeasured = 0
    for idx, tag in enumerate(tags, start=1):
        ent = nm.get(tag)
        if not ent or not ent[0]:      continue
        sid = ent[0]
        if sid not in draft:           continue
        if idx not in notes_xml:
            unmeasured += 1; continue      # 🟥 조용히 넘기면 «대조 N절» 의 N 만 줄어 «차이 0» 쪽으로 열화된다
        if sid in skip_units:
            skipped += 1; continue
        a = pn.norm_lines(draft[sid]); b = pn.norm_lines(notes_xml[idx])
        compared += 1
        if a != b:
            # 🟥 «다르다» 를 한 덩어리로 내면 처방이 뭉개진다. 두 부류는 처방이 정반대다:
            #    ORDER-DRIFT = 재빌드(내용은 다 있다) · GETPUT-LOSS = 왕복기 수리(내용이 없다)
            owner = next((o for o in draft
                          if o != sid and o not in skip_units and pn.norm_lines(draft[o]) == b), None)
            if owner:
                findings.append(('roundtrip', 'ORDER-DRIFT', f'slide{idx}',
                                 f'배선은 {sid} 인데 산출물엔 {owner} 의 낭독이 들어 있다',
                                 '내용 손실 아님 — 산출물이 현재 배선 순서를 안 따른다(재빌드 필요). '
                                 '🟥 그 전에 그 산출물에 사람 손첨삭이 있는지 봐라 — 재빌드가 덮는다'))
            else:
                only = len([x for x in a if x not in b]) + len([x for x in b if x not in a])
                findings.append(('roundtrip', 'GETPUT-LOSS', sid,
                                 f'{len(a)}문장 vs {len(b)}문장 · 불일치 {only}',
                                 '사람 편집 0 구간인데 왕복이 안 닫힌다 — 차이는 왕복기 손실이다'))
    notes.append(f'L4 getput: 대조 {compared}절 · 근사구간 skip {skipped}절 (S5·S15, 왕복기 명시 한계) '
                 f'· 🟥 UNMEASURED {unmeasured}절 (노트 매핑 없음 — 0 아님)')
    if no_note:
        notes.append(f'L4 노트 없는 슬라이드 {len(no_note)}장 (위치 {no_note[:8]}{"…" if len(no_note)>8 else ""}) — UNMEASURED')
    if compared == 0:
        notes.append('L4 🟥 UNMEASURED — 대조된 절이 0. 컨트롤이 죽었다, «차이 0» 아님')
    return findings, notes

# ── main ──────────────────────────────────────────────────────────────────────
# ── L5 jargon ─────────────────────────────────────────────────────────────────
def lane_jargon(cfg, root, texts, surf_meta):
    """선언된 조어가 **낭독면 첫 등장에 풀이를 달고 있는가**.

    🟥 판정하는 것은 «이 단어가 어려운가»가 아니다 — 그건 결론이고 코드로 굳히면 천장이 된다.
       재는 것은 기록의 성질 하나다: 선언 목록에 있는 어휘의 첫 USE 둘레에 풀이 어휘가 있는가.
       어느 단어를 조어로 볼지는 `jargon_terms.yaml` 에 사람이 적는다.
    🟥 낭독면(spoken: true)만 본다. 화면 텍스트·원장·Q&A 는 다른 규약이라 여기 섞으면 오탐이다.
    미선언 어휘는 구조적으로 안 보인다 — 이 레인의 0 은 «조어가 없다»가 아니라
    «선언한 것 중 풀이 없는 게 없다»이다. 0 을 «깨끗하다»로 읽지 말 것."""
    cfgp = cfg.get('jargon_terms')
    if not cfgp: return [], ['L5 jargon : 선언 파일 없음 — NOT_CONFIGURED (0 아님)']
    # 🟥 2026-09-04 — canon_terms 와 같은 수리(신호 §3-2·3-3): root 우선 · HERE 폴백 ·
    #    부재는 크래시가 아니라 UNMEASURED.
    jt_path, jt_fb = resolve_declared(root, cfgp)
    jt_notes = []
    if jt_fb:
        jt_notes.append(f'L5 jargon : jargon_terms 를 root 에서 못 찾아 스킬 위치로 폴백 '
                        f'(skill-dir fallback) — {jt_path}')
    try:
        spec = load(jt_path)
    except Exception as e:
        return [], [f'L5 jargon : jargon_terms 부재({jt_path}) — UNMEASURED (0 아님): '
                    f'{type(e).__name__}: {e}']
    win = int(spec.get('gloss_window', 2))
    findings, notes = [], list(jt_notes)
    spoken = [sid for sid, m in surf_meta.items() if m.get('spoken')]
    if not spoken:
        return [], notes + ['L5 jargon : 낭독면 선언 0 — UNMEASURED (0 아님)']
    checked = glossed = 0
    for sid in spoken:
        if sid not in texts: continue
        text, _ = texts[sid]
        # 🟥 대상이 스스로 «STALE / 고치지 마라» 라고 선언했으면 판정 대상이 아니다.
        # L2 가 이미 «대상의 은퇴 선언 어휘를 존중»하는 것과 같은 원리 — 안 하면 폐기된
        # 사본을 고치라고 지시하게 되고, 그건 이 코퍼스가 이미 세 번 겪은 이중 소스 사고다.
        head = '\n'.join(text.split('\n')[:40])
        if 'STATUS: STALE' in head or '이 파일을 고치지 마라' in head:
            notes.append(f'L5 jargon : {sid} 는 자기선언 STALE — 판정 제외 (0 아님, RETIRED)')
            continue
        lines = text.split('\n')
        marker = surf_meta[sid].get('marker')
        for t in spec['terms']:
            first = None
            for lit in t['literals']:
                for cls, no, ln in hits(text, lit, marker):
                    if cls != 'USE': continue
                    # 🟥 이력/주석 블록 배제 (2026-08-20 첫 실사용에서 이 레인이 낸 오탐).
                    # 이 코퍼스의 개정 이력은 낭독과 **같은 `>` 인용 블록**에 살고, 앞머리
                    # 이모지로만 갈린다. 전역 classify 를 고치지 않는 이유: L1 이 그 함수에
                    # 붙어 있고 이 델타는 L5 만 고치기로 한 범위다(회귀 위험 > 중복 비용).
                    if _HIST.match(ln): continue
                    if first is None or no < first[0]: first = (no, lit, ln)
            if first is None: continue
            checked += 1
            no, lit, ln = first
            lo, hi = max(0, no - 1 - win), min(len(lines), no + win)
            around = '\n'.join(lines[lo:hi])
            if any(g in around for g in t.get('gloss', [])):
                glossed += 1
                findings.append((f'{sid}:L5', 'mention', f'{t["id"]}@{no}', lit,
                                 '첫 등장에 풀이 있음 — 통과'))
            else:
                findings.append((sid, 'JARGON-UNGLOSSED', f'{t["id"]}@{no}', lit,
                                 f'선언된 조어인데 첫 등장 ±{win}줄에 풀이가 없다 '
                                 f'▸ 권장: {t["plain"]} ▸ {ln[:60]}'))
    notes.append(f'L5 jargon : 선언 {len(spec["terms"])}항 · 낭독면 등장 {checked}항 · '
                 f'첫 등장 풀이 있음 {glossed}항 '
                 f'(⚠️ 미선언 어휘는 구조적으로 안 보인다 — 0 을 «조어 없음»으로 읽지 말 것)')
    return findings, notes


def lane_pacing(cfg, root, texts, surf_meta):
    """L6 — **장당 발화가 상한을 넘는가.** 넘으면 쪼개거나 줄이라고 «권고»한다.

    🟥 재는 것은 기록의 성질 하나다: 이 장에 붙은 낭독의 길이. «이 장이 지루한가»가 아니다
       — 그건 결론이고, 코드로 굳히면 오늘의 취향이 내일의 천장이 된다.
    🟥 상한은 **`surfaces.yaml speech_pacing` 에 사람이 적는다.** 문헌엔 «한 장 N초»가 없다
       (그 근거는 그 yaml 주석에 적어 뒀다). 기본값은 이 코퍼스 n=1 실측이므로 새 발표는
       분포부터 다시 재라 — 그래서 이 레인은 상한 위반뿐 아니라 **평균·중앙값·이봉 여부**를
       항상 같이 출력한다. 분포를 안 보여주면 남의 숫자를 그대로 쓰게 된다.
    🟥 **단위가 «장»인지 «절»인지 반드시 밝힌다.** 빌드된 덱(pptx)이 있으면 장 단위로 재고,
       없으면 원고의 절 단위로 재면서 그 사실을 라벨에 박는다 — 한 절이 여러 장으로 쪼개져
       있으면 절 단위 수치는 **과대**다. 조용히 절을 장이라고 부르면 그 순간 계기가 거짓말한다.
    """
    pac = cfg.get('speech_pacing')
    if not pac:
        return [], ['L6 pacing : speech_pacing 선언 없음 — NOT_CONFIGURED (0 아님)']
    lo = float(pac.get('chars_per_sec', [5.18, 5.35])[0])
    hard, soft = float(pac.get('hard_sec', 45)), float(pac.get('soft_sec', 30))

    unit, rows = None, []
    deck = cfg.get('surfaces_by_id', {}).get('built_deck')
    dpath = resolve(root, deck['path']) if deck else None
    if dpath and os.path.exists(dpath):
        try:
            from pptx import Presentation                       # noqa
            prs = Presentation(dpath)
            for i, sl in enumerate(prs.slides, 1):
                t = sl.notes_slide.notes_text_frame.text
                if t.startswith('[간지]') or t.startswith('[진행]') or not t.strip():
                    continue
                rows.append((f'p{i}', len(re.sub(r'\s', '', t.replace('(이어서) ', '')))))
            unit = '장'
        except ImportError:
            pass
    if unit is None:
        # 🟥 절 단위 폴백. **첫 구현이 6.5배 틀렸다**(18.0분짜리 원고를 116.8분으로 셌다) —
        #   낭독과 개정 이력이 **같은 `>` 인용 블록**에 살아서, 「> 로 시작하면 낭독」 규칙이
        #   이력을 통째로 삼켰다. 6.5배 틀린 계기에 «절 단위라 과대» 라벨만 달아 내보내면
        #   그 라벨이 오히려 오차를 정상으로 보이게 만든다 — 라벨은 오차의 변명이 아니다.
        #   ⇒ 규칙을 정확히 옮긴다: 🗣 뒤의 `>` 를 **빈 줄로 블록 분할**하고, 각 블록의
        #     **첫 글자가 한글/영숫자인 것만** 낭독으로 센다(이력 블록은 🆕/🟥/⚠️ 로 시작한다).
        spoken = [sid for sid, m in surf_meta.items() if m.get('spoken')]
        if not spoken:
            return [], ['L6 pacing : 낭독면 선언 0 — UNMEASURED (0 아님)']
        _STRIP = '*_"\'«»『「(> '
        for sid in spoken:
            if sid not in texts: continue
            text, _ = texts[sid]
            marker = surf_meta[sid].get('marker') or '\U0001f5e3'
            unit_re = re.compile(surf_meta[sid].get('unit_pattern') or r'^###\s+(\S+)')
            state = {'cur': None, 'blocks': [], 'blk': []}
            def _flush():
                if state['blk']:
                    state['blocks'].append(list(state['blk'])); state['blk'].clear()
            def _emit():
                if not state['cur']: return
                body = []
                for b in state['blocks']:
                    probe = b[0].lstrip(_STRIP)
                    if not probe: continue
                    c = probe[0]
                    if not (('가' <= c <= '힣') or c.isalnum()): continue
                    body.extend(b)
                joined = re.sub(r'\*\*|`|\*', '', ''.join(body))
                rows.append((state['cur'], len(re.sub(r'\s', '', joined))))
            mode = 'out'
            for raw in text.split('\n'):
                st = raw.strip()
                m = unit_re.match(raw)
                if m:
                    _flush(); _emit()
                    state['cur'], state['blocks'], mode = m.group(1), [], 'out'; continue
                if state['cur'] is None: continue
                if st.startswith(marker): _flush(); mode = 'speech'; continue
                if st[:1] in ('\U0001f5a5', '⏱') or st.startswith('##'):
                    _flush(); mode = 'out'; continue
                if mode != 'speech': continue
                if st == '' or st == '---': _flush(); continue
                if not st.startswith('>'): _flush(); mode = 'out'; continue
                state['blk'].append(st.lstrip('> '))
            _flush(); _emit()
        unit = '절'
    # 🟥 배선에 없는 절은 **무대에 안 선다** — 세면 과대다. 실측: 폴백이 26절 21.2분을 냈는데
    #   배선된 절만 세면 24절이고, 차이 2절이 그대로 3.2분이었다. 배선 표면이 있으면 거른다.
    if unit == '절':
        wsurf = cfg.get('surfaces_by_id', {}).get('wiring')
        try:
            live = note_map_ids(resolve(root, wsurf['path'])) if wsurf else None
        except Exception as e:
            live = None; notes_pre = f'배선 파싱 실패({e})'
        else:
            notes_pre = None
        if live:
            dead = [k for k, _ in rows if k not in live]
            rows = [(k, n) for k, n in rows if k in live]
            if dead: _dead_note = f'L6 pacing : 배선에 없어 제외한 절 {len(dead)}개 — {", ".join(dead)}'
            else:    _dead_note = 'L6 pacing : 배선에 없는 절 0'
        else:
            _dead_note = ('L6 pacing : 🟥 배선 대조 못 함' + (f' ({notes_pre})' if notes_pre else '')
                          + ' — 무대에 안 서는 절이 섞여 **과대**일 수 있다. 이건 0 이 아니라 미대조다')
    else:
        _dead_note = None
    rows = [(k, n) for k, n in rows if n]
    if not rows:
        return [], ['L6 pacing : 낭독 추출 0 — UNMEASURED (0 아님, 계기가 못 읽은 것일 수 있다)']

    secs = sorted(n / lo for _, n in rows)
    mean, med = sum(secs) / len(secs), secs[len(secs) // 2]
    # 이봉 판정 — 상한 근처 ±5초 대역이 전체의 5% 미만이면 «사이가 비었다»
    band = [x for x in secs if abs(x - soft) <= 5]
    findings, notes = [], []
    if _dead_note: notes.append(_dead_note)
    notes.append(f'L6 pacing : 단위={unit} · {len(rows)}개 · 평균 {mean:.0f}초 · 중앙값 {med:.0f}초 '
                 f'· 총 {sum(secs)/60:.1f}분 (자수÷실측속도 — 전환·호흡 제외. '
                 f'**실낭독 실측·배분 합은 원고 ⏱ 가 정본이고 이 값보다 길다**) '
                 f'· 상한 HARD {hard:.0f}/SOFT {soft:.0f}초'
                 + ('' if unit == '장' else '  ⚠️ 절 단위라 장 단위보다 **과대**하다'))
    notes.append(f'L6 pacing : SOFT±5초 대역에 {len(band)}개({len(band)/len(rows)*100:.0f}%) — '
                 + ('분포가 갈려 있다(이봉). 상한을 몇 초로 두든 걸리는 것이 거의 같다'
                    if len(band)/len(rows) < 0.05 else
                    '경계 대역에 몰려 있다 — 상한을 바꾸면 결과가 크게 흔들린다. 숫자를 신중히'))
    for k, n in sorted(rows, key=lambda r: -r[1]):
        sec = n / lo
        if sec > hard:
            findings.append((k, 'PACING-OVER-HARD', f'{n}자', f'{sec:.0f}초',
                             f'{unit} 하나가 상한 {hard:.0f}초를 넘는다 ▸ 쪼개거나 줄여라'))
        elif sec > soft:
            findings.append((f'{k}:L6', 'mention', f'{n}자', f'{sec:.0f}초',
                             f'권고 대역 초과({soft:.0f}초) — 판단은 사람'))
    return findings, notes


# ── L8 interslide ─────────────────────────────────────────────────────────────
def lane_interslide(cfg, root):
    """L8 — **장 사이 의존이 깨졌나.** 챔버 런 #15 가 배출한 계기를 하네스 안에서 돌린다.

    실사고: S8 을 빼자 S9 가 성립 불가가 됐다(S9 는 S8 의 화면을 재사용). 렌더를 한 장씩 다 봐도 구조적으로 안 잡혔다 — **운영자가 잡았다** — 한 장씩 보기 때문이다. 사람이 잡았지 계기가 아니었다.

    🟥 **advisory 고정 — findings 를 내지 않는다(종료코드에 안 태운다).**
       판정문(챔버 #15 EMISSION_VERDICT)이 못박은 운용 형태다. 이 계기의 오탐을
       최저비용으로 무마하는 길이 «지시어 삭제»라서, 종료코드에 태우면 원고를
       **나쁜 방향으로 미는** 압력이 된다.
    ⚠️ UNRESOLVED 는 «안 깨졌다»가 아니라 **미측정**이다. 0 으로 접지 않는다.
    """
    spec = cfg.get('interslide')
    if not spec:
        return [], ['L8 interslide : 선언 없음 — NOT_CONFIGURED (0 아님)']
    mod_p = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'interslide_deps.py')
    if not os.path.exists(mod_p):
        return [], [f'L8 interslide : 계기 부재({mod_p}) — NOT_WIRED (0 아님)']
    surf_by_id = {x['id']: x for x in cfg.get('surfaces', [])}
    ms = surf_by_id.get(spec.get('manuscript'))
    if not ms:
        return [], [f'L8 interslide : 원고 표면 «{spec.get("manuscript")}» 미선언 — NOT_CONFIGURED (0 아님)']
    draft_p, build_p = resolve(root, ms['path']), resolve(root, spec['builder'])
    for label, p in (('원고', draft_p), ('배선 정본', build_p)):
        if not os.path.exists(p):
            return [], [f'L8 interslide : {label} 실물 없음({p}) — UNMEASURED (0 아님)']
    try:
        import importlib.util
        sp = importlib.util.spec_from_file_location('interslide_deps', mod_p)
        m = importlib.util.module_from_spec(sp); sp.loader.exec_module(m)
        # 🟥 2026-09-04 — 원고 표면의 unit_pattern 선언을 analyze() 로 그대로 넘긴다
        #    (신호 §3-1 수리 — units() 하드코딩을 없앤 절반이 여기다).
        r = m.analyze(draft_p, build_p, pattern=ms.get('unit_pattern'))
    except Exception as e:
        return [], [f'L8 interslide : 계기 오류({e}) — UNMEASURED (0 아님)']
    if r.get('error'):
        return [], [f'L8 interslide : {r["error"]} — UNMEASURED (0 아님)']

    hard = sum(1 for e in r['edges'] if e[2] == 'HARD')
    notes = [f'L8 interslide : 절 {r["n_units"]} · 배선 {r["n_wired"]} · HARD 의존 {hard}건 '
             f'· 깨짐 {len(r["broken"])} · UNRESOLVED {len(r["unresolved"])} (advisory)']
    for sid, prev, _ in r['broken']:
        notes.append(f'   {ANSI_BAD} L8 깨진 의존: {sid} 가 {prev} 를 필요로 하는데 {prev} 는 안 만들어진다')
    for sid, terms, _ in r['unresolved']:
        notes.append(f'   {ANSI_UNK} L8 UNRESOLVED: {sid} — 대상어 {terms} (미측정, 0 아님)')
    return [], notes

def lane_regen(cfg, root):
    """L7 — **재생성이 사람 편집을 삼켰나.** 생성면 ↔ 사람이 만진 사본을 대조한다.

    🟥 이 레인은 README 「다음(미착수)」 2번(**표면별 정본 축 선언** — 문자열=원고 /
       서식=Keynote, Jupytext paired-notebook 원리)이 **실제로 터진 뒤에** 생겼다.
       2026-08-21 실측: 빌더가 발표자 노트를 **전량 덮어써서** 사람이 손으로 쓴 발화
       **20곳**이 소리 없이 사라졌다. 빌드는 성공했고, 원고↔빌드 대조는 «원고대로»였다.
       ⇒ 대조 축이 «원본 ↔ 생성물» 하나뿐이면 **사람 편집면은 구조적으로 안 보인다.**

    🟥 최소 길이 필터를 두지 않는다(minlen 없음). 같은 날 실측: 회수 도구에 12자 필터를
       걸었더니 「계측기→검증 장치」·「—→.」 같은 **짧지만 의도가 분명한 수정 14곳**이
       통째로 안 보였다. 필터는 노이즈를 줄인 게 아니라 **신호를 지웠다.**

    ⚠️ 이 레인의 «0» 은 «두 사본이 같다»가 아니라 **«사람 사본에만 있는 조각이 없다»**이다.
       생성물에만 있는 것(원고에서 온 정상 갱신)은 세지 않는다 — 방향이 있는 대조다.
    """
    spec = cfg.get('regen_recovery')
    if not spec:
        return [], ['L7 regen : regen_recovery 선언 없음 — NOT_CONFIGURED (0 아님)']
    gen_p, hum_p = resolve(root, spec['generated']), resolve(root, spec['human'])
    for label, p in (('generated', gen_p), ('human', hum_p)):
        if not os.path.exists(p):
            return [], [f'L7 regen : {label} 실물 없음({p}) — UNMEASURED (0 아님). '
                        f'사람 사본은 **내보내야** 생긴다(키노트 → pptx). 안 내보내면 이 레인은 못 잰다']
    try:
        from pptx import Presentation
    except ImportError:
        return [], ['L7 regen : python-pptx 없음 — UNMEASURED (0 아님)']
    import difflib

    def notes(path):
        pr = Presentation(path)
        return [re.sub(r'\s+', ' ', (s.notes_slide.notes_text_frame.text
                                     if s.has_notes_slide else '')
                       .replace('(이어서) ', '')).strip() for s in pr.slides]
    A, B = notes(gen_p), notes(hum_p)
    findings, notes_out = [], []
    if len(A) != len(B):
        notes_out.append(f'L7 regen : 🟥 장 수가 다르다 — 생성 {len(A)} · 사람 {len(B)}. '
                         f'앞에서부터 짝지으므로 **어긋난 뒤 결과는 못 믿는다**(부분 측정)')
    for i, (a, b) in enumerate(zip(A, B), 1):
        if a == b:
            continue
        sm = difflib.SequenceMatcher(None, a, b)
        gone = [b[j1:j2] for t, _, _, j1, j2 in sm.get_opcodes()
                if t in ('replace', 'insert') and b[j1:j2].strip()]
        for g in gone:
            findings.append((f'p{i}', 'REGEN-ATE-EDIT', '사람 사본에만 있음', g.strip()[:70],
                             '재생성이 삼켰다 — **원본(생성 소스)으로 옮겨라.** 사본에 두면 다음 빌드가 또 지운다'))
    notes_out.append(f'L7 regen : {min(len(A), len(B))}장 대조 · 사람 사본에만 있는 조각 {len(findings)}건 '
                     f'(⚠️ 0 은 «같다»가 아니라 «삼킨 게 없다»이다)')
    return findings, notes_out


def selftest_pacing():
    """known-pair. 계기가 «상한 초과»와 «이내»를 실제로 가르는지 — 안 가르면 장식이다."""
    cfg = {'speech_pacing': {'chars_per_sec': [5.0, 5.0], 'hard_sec': 10, 'soft_sec': 5}}
    meta = {'s': {'spoken': True, 'marker': '🗣', 'unit_pattern': r'^###\s+(\S+)'}}
    long_ = '가' * 60      # 60자 / 5.0 = 12초 > HARD 10
    short = '가' * 20      # 20자 / 5.0 =  4초 < SOFT 5
    cases = [
        ('known-positive(상한 초과)', f'### A\n🗣\n> {long_}\n', 1),
        ('known-negative(상한 이내)', f'### A\n🗣\n> {short}\n', 0),
        ('control(낭독 마커 없음)',   f'### A\n> {long_}\n',      0),
    ]
    ok = True
    for name, body, want in cases:
        f, _ = lane_pacing(cfg, '.', {'s': (body, 'markdown')}, meta)
        got = len([x for x in f if x[1] == 'PACING-OVER-HARD'])
        hit = got == want
        ok &= hit
        print(f'  L6 {name}: got {got} want {want}  {"PASS" if hit else "FAIL"}')
    print('L6 SELFTEST:', 'PASS' if ok else 'FAIL — 이 레인의 출력을 쓰지 마라')
    return 0 if ok else 1


def selftest_jargon():
    """known-pair. 계기가 «풀이 있음»과 «없음»을 실제로 가르는지 — 안 가르면 이 레인은 장식이다."""
    import tempfile, textwrap
    ok = True
    spec = {'gloss_window': 2, 'terms': [
        {'id': 'T', 'literals': ['탈상관'], 'plain': 'p', 'gloss': ['만든 회사'], 'why': 'w'}]}
    cases = [('known-negative(풀이 있음)', '🗣\n> 셋째, 탈상관입니다.\n> 만든 회사를 섞습니다.\n', 0),
             ('known-positive(풀이 없음)', '🗣\n> 셋째, 탈상관입니다.\n> 그래서 섞습니다.\n', 1),
             ('control(어휘 부재)',        '🗣\n> 아무 상관 없는 문장입니다.\n', 0)]
    with tempfile.TemporaryDirectory() as d:
        sp = os.path.join(d, 'j.yaml')
        with open(sp, 'w', encoding='utf-8') as f: yaml.safe_dump(spec, f, allow_unicode=True)
        for label, body, expect in cases:
            f, _ = lane_jargon({'jargon_terms': sp}, d, {'s': (body, 'markdown')},
                               {'s': {'marker': '🗣', 'spoken': True}})
            got = len([x for x in f if x[1] == 'JARGON-UNGLOSSED'])
            good = got == expect
            ok &= good
            print(f"  {'PASS' if good else 'FAIL'}  {label}: expect={expect} got={got}")
    print('L5 SELFTEST:', 'PASS — 계기가 짝을 가른다' if ok else 'FAIL — 판별력 없음')
    return 0 if ok else 2


def main():
    if '--self-test' in sys.argv: return selftest_jargon() or selftest_pacing()
    # 🟥 2026-09-01 — «자기 판»을 config 로드 «전»에 찍는다. 오늘 실사고: 갈라진 옛 포크
    #    (동반 저장소의 preprep, 724줄)를 돌려놓고 「설치본에 레인이 없다」고 보고했다 — 레인은
    #    정본(776줄)에 있었고 귀속이 틀렸다. 🟥 첫 판은 이 줄을 `cfg` 로드 «뒤»에 뒀는데,
    #    config 가 없으면 안 찍힌다 — 「내가 어느 판인가」가 «가장 필요한 순간»에 침묵한다.
    #    ⚠️ 내장 `lane_*` 함수만으로는 판별이 «안 된다»(정본·포크 둘 다 7). 갈리는 것은
    #    별도 모듈 레인(L9/L10/L11)이라 둘 다 센다.
    #    [[feedback_rule_misdescribes_its_own_machine]] · [[feedback_instrument_vs_target_and_budget]]
    _lanes = sorted(n[5:] for n in globals() if n.startswith('lane_') and callable(globals()[n]))
    _mods = sorted(m for m in ('lane_promise', 'lane_adjacent_dup', 'lane_progression',
                                'lane_slide_relations', 'lane_geometry', 'lane_diagram', 'lane_slide_refs')
                   if os.path.exists(os.path.join(HERE, m + '.py')))
    # 🟥 2026-09-04 신설 — `--lane R1,R2,...` 로 **새 모듈 레인(R1-R5·P1/P3)만** 골라 끈다/켠다.
    #    L1~L11 은 아직 이 필터를 안 탄다(usage 줄의 --lane 은 그쪽엔 미배선인 채 남아 있다 —
    #    있는 척하지 않는다). 최소 접촉 원칙: 기존 L* 배선은 안 건드린다.
    _lane_arg = sys.argv[sys.argv.index('--lane') + 1] if '--lane' in sys.argv else None
    _wanted = set(_lane_arg.split(',')) if _lane_arg else None

    def _lane_on(*codes):
        return _wanted is None or any(c in _wanted for c in codes)
    print(f"── 발표 준비 하네스 v0.1 ──")
    print(f"   판: {os.path.realpath(__file__)}")
    print(f"   레인 {len(_lanes)}(내장): {' · '.join(_lanes)}")
    print(f"   레인 {len(_mods)}(모듈): {' · '.join(m[5:] for m in _mods) if _mods else '🟥 없음 — 갈라진 사본일 수 있다'}")
    cfg_path = sys.argv[1] if len(sys.argv) > 1 and not sys.argv[1].startswith('--') \
               else os.path.join(HERE, 'surfaces.yaml')
    cfg = load(cfg_path)
    root = os.path.expanduser(cfg['root'])
    surf = {s['id']: s for s in cfg['surfaces']}


    # L3 inventory — 먼저 돈다. 무엇을 못 읽었는지가 뒤 레인의 계상 범위를 정한다
    texts, missing, unreadable = {}, [], []
    for s in cfg['surfaces']:
        p = resolve(root, s['path'])
        if not os.path.exists(p): missing.append(s['id']); continue
        try: texts[s['id']] = (read_text(p, s['kind']), s['kind'])
        except Exception as e: unreadable.append((s['id'], str(e)))
    print(f"L3 inventory : 읽음 {len(texts)} · 부재 {len(missing)} · 기계로 못 읽음 {len(unreadable)}")
    for m in missing:     print(f"   {ANSI_BAD} 선언됐는데 디스크에 없다: {m}  (UNMEASURED, 0 아님)")
    for i, why in unreadable: print(f"   {ANSI_UNK} 기계 어댑터 없음: {i} — {why}  (UNMEASURED)")

    findings, notes = [], []
    # 🟥 2026-09-04 — 'unit_pattern' 을 추가했다. L6 pacing 은 이미 이 필드를 읽고 있었는데
    #    (surf_meta[sid].get('unit_pattern')) 여기서 채워준 적이 없어 **항상 None** 이었다 —
    #    L10/L11 도 같은 필드로 units() 에 pattern 을 전달하므로 여기가 단일 공급처가 된다.
    surf_meta = {s['id']: {'marker': s.get('speech_marker'), 'spoken': s.get('spoken', False),
                           'unit_pattern': s.get('unit_pattern')}
                 for s in cfg['surfaces']}
    f1, n1 = lane_canon(cfg, root, texts, surf_meta); findings += f1; notes += n1
    try:
        import lane_promise
        import importlib.util as _iu2
        _sp2 = _iu2.spec_from_file_location('interslide_deps', os.path.join(HERE, 'interslide_deps.py'))
        _im2 = _iu2.module_from_spec(_sp2); _sp2.loader.exec_module(_im2)
        _f11, _n11 = lane_promise.scan(texts, surf_meta, mod=_im2)
        notes += _n11   # advisory 고정 — findings 에 안 태운다
    except Exception as _e:
        notes.append('L11 promise : 계기 미실행 — NOT_WIRED (%s: %s) (0 아님)' % (type(_e).__name__, _e))
    try:
        import lane_adjacent_dup, importlib.util as _ilu
        _mp = os.path.join(HERE, 'interslide_deps.py')
        _sp = _ilu.spec_from_file_location('interslide_deps', _mp)
        _im = _ilu.module_from_spec(_sp); _sp.loader.exec_module(_im)
        _f10, _n10 = lane_adjacent_dup.scan(cfg, texts, surf_meta, _im)
        findings += _f10; notes += _n10
    except Exception as _e:
        notes.append('L10 adjacent-dup : 계기 미실행 — NOT_WIRED (%s: %s) (0 아님)' % (type(_e).__name__, _e))
    try:
        import lane_progression
        _f9, _n9 = lane_progression.scan(cfg, texts, surf_meta)
        findings += _f9; notes += _n9   # 차단 — 선언된 것만 보므로 오탐이 선언에 갇혀 있다
    except Exception as _e:
        # 계기 부재/오류를 «통과»로 렌더하지 않는다. 부재 != 0.
        notes.append('L9 progression : 계기 미실행 — NOT_WIRED (%s: %s) (0 아님)' % (type(_e).__name__, _e))

    f2, n2 = lane_unitref(cfg, root, surf); findings += f2; notes += n2
    f4, n4 = lane_getput(cfg, root, surf); findings += f4; notes += n4
    f5, n5 = lane_jargon(cfg, root, texts, surf_meta); findings += f5; notes += n5
    # L6 은 built_deck 표면을 장 단위로 읽고 싶어 하므로 id→선언 맵을 넘겨준다
    cfg['surfaces_by_id'] = {x['id']: x for x in cfg.get('surfaces', [])}
    f6, n6 = lane_pacing(cfg, root, texts, surf_meta); findings += f6; notes += n6
    f7, n7 = lane_regen(cfg, root); findings += f7; notes += n7
    f8, n8 = lane_interslide(cfg, root); findings += f8; notes += n8
    if _lane_on('R1', 'R2', 'R3', 'R4', 'R5', 'R'):
        try:
            import lane_slide_relations
            _fR, _nR = lane_slide_relations.scan(cfg, root)
            notes += _nR   # advisory 고정 — findings 에 안 태운다(L8·L11 관례)
        except Exception as _e:
            notes.append('R1-R5 slide-relations : 계기 미실행 — NOT_WIRED (%s: %s) (0 아님)'
                         % (type(_e).__name__, _e))
    # L13 slide-refs — 「N p」 리터럴 장 번호 참조가 실재하는 장을 가리키나(2026-09-06 신설).
    #   이웃 레인들은 «말로 가리키는 것»만 본다(L5/L8 화면 재사용·시각 지시 · R1~R5 접속어).
    #   숫자 참조는 아무도 안 봤고, 장을 빼면 그 뒤가 통째로 당겨진다.
    #   차단은 «범위 밖» 하나뿐 — 「그 번호가 맞는 장인가」는 의미 판단이라 안 본다.
    try:
        import lane_slide_refs
        _f13, _n13 = lane_slide_refs.scan(cfg, root)
        findings += _f13; notes += _n13
    except Exception as _e:
        notes.append('L13 slide-refs : 계기 미실행 — NOT_WIRED (%s: %s) (0 아님)'
                     % (type(_e).__name__, _e))
    # L12 diagram — 타입 JSON 도해가 «지금 JSON» 에서 validate 를 거쳐 구워졌나(2026-09-05 신설).
    #   차단: 선언된 diagram_source 표면에 한해 지문·validate·해상도·viewBox 폭·여백을 본다.
    #   영수증 없는 PNG(손그림)는 UNMEASURED 노트로만 — 0 이 아니다.
    try:
        import lane_diagram
        _f12, _n12 = lane_diagram.scan(cfg, root)
        findings += _f12; notes += _n12
    except Exception as _e:
        notes.append('L12 diagram : 계기 미실행 — NOT_WIRED (%s: %s) (0 아님)' % (type(_e).__name__, _e))
    if _lane_on('P1', 'P3', 'P'):
        try:
            import lane_geometry
            _fP, _nP = lane_geometry.scan(cfg, root)
            notes += _nP   # advisory 고정 — findings 에 안 태운다(L8·L11 관례)
        except Exception as _e:
            notes.append('P1/P3 geometry : 계기 미실행 — NOT_WIRED (%s: %s) (0 아님)'
                         % (type(_e).__name__, _e))

    for n in notes: print(f"   · {n}")
    use = [f for f in findings if f[1] != 'mention']
    men = [f for f in findings if f[1] == 'mention']
    print(f"\n판정 대상(USE) {len(use)}건 · 참고(MENTION — «쓰지 마라»고 적은 줄) {len(men)}건")
    print("🟥 둘을 합치지 않는다. MENTION 은 0 이 아니라 **다른 것**이다.")
    for sid, kind, ref, lit, why in use:
        print(f"  {ANSI_BAD} [{kind}] {sid} ← {ref}: «{lit}»\n       {why}")
    if men:
        print(f"\n  ({ANSI_UNK} MENTION {len(men)}건 — 표면별: "
              + ', '.join(f'{s}×{sum(1 for m in men if m[0].split(":")[0]==s)}'
                          for s in sorted({m[0].split(':')[0] for m in men})) + ")")
    findings = use

    print(f"\n{ANSI_UNK} 계기 한계(명시): pptx 는 XML 텍스트 레이어만 본다 — "
          f"이미지에 구워진 텍스트는 구조적으로 안 보인다. Keynote 는 사람만 읽는다.")
    if unreadable or missing: return 2 if not findings else 1
    return 1 if findings else 0

# 🟥 `__main__` 가드 — 이게 없으면 **import 가 곧 실행**이라 아무도 이 파일을 라이브러리로
#    못 쓴다(레인 하나만 부르려 해도 전체가 돌고 종료한다). 배포본이 된 뒤에 드러난 결함이다:
#    스킬 문서가 «레인을 골라 부를 수 있다»는 인상을 주는데 실제로는 CLI 한 형태뿐이었다.
if __name__ == '__main__':
    sys.exit(main())
