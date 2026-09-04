#!/usr/bin/env python3
"""pptx-only 자산군 어댑터 — preprep 레인이 읽을 수 있는 형태를 pptx 에서 세운다.

preprep 은 「원고 마크다운 + 빌더 스크립트」를 전제한다. **원고·빌더가 없고 pptx 자체가
정본**인 자산군엔 그동안 붙일 길이 없었다(신호 §3-4). 이 어댑터가 그 자리를 메운다 — pptx 에서
세 가지를 뽑아 세운다:

  <out>/manuscript.md   장마다 🖥(화면 텍스트) · 🗣(발표자 노트). 둘 다 pptx 에서 뽑는다
  <out>/note_map.py     NOTE_MAP — 어느 절이 실제 장으로 배선됐나 (이 어댑터는 절=장 1:1)
  <out>/surfaces.yaml   표면 선언. canon_terms/jargon_terms 는 이 파일의 basename 으로만
                        적는다 — root(=<out>) 기준으로 못 찾으면 이제 스킬 위치로 자동
                        폴백하므로(2026-09-04 수리), 사람이 자기 목록을 만들기 전엔 배포본의
                        example 을 그대로 가리켜도 UNMEASURED 로 안 죽는다

🟥 이건 «대본을 다시 만드는 것»이 아니다 — **읽기 전용 파생물**이고, 원본 pptx 는 절대 안
   건드린다. 손으로 고치지 마라 — 재실행이 덮어쓴다.
🟥 NOTE_MAP 이 전 절을 1:1 로 담으므로 L8 의 BROKEN(배선 끊김)은 이 어댑터 산출물에서
   구조적으로 안 뜬다. 값을 내는 것은 UNRESOLVED — «화면/앞 절을 가리키는데 대상을 못
   찾겠다» — 뿐이다.

🟥 **함정(신호 §3-4 실측, 재발 방지로 이 주석과 known-pair 를 같이 둔다)**: 초판은 「첫
   도형은 제목이니 빼고 나머지만 🖥 로 담자」로 짰다가, **제목 도형 하나뿐인 목차형 장의
   화면이 통째로 비었다.** 첫 도형도 화면이다 — `_shapes()` 는 절대 첫 항목을 자르지 않는다.
   `--self-test` 가 known-positive(제목 하나뿐인 장)로 이 재발을 계속 감시한다.

usage:
  preprep_wire.py <deck.pptx> [--out DIR]
  preprep_wire.py --self-test
"""
import zipfile, re, os, sys

# 텍스트가 있어도 화면으로 안 셀 도형의 name(예: 장식용 섹션 안내 도형). 프로젝트마다
# 다르므로 사람이 편집하는 자리다 — 비워 두면 아무것도 안 거른다.
SKIP_SHAPES = ()


def _order(z):
    rels = dict(re.findall(r'Id="([^"]+)"[^>]*Target="slides/slide(\d+)\.xml"',
                           z.read('ppt/_rels/presentation.xml.rels').decode('utf-8')))
    lst = re.search(r'<p:sldIdLst>.*?</p:sldIdLst>',
                    z.read('ppt/presentation.xml').decode('utf-8'), re.S).group(0)
    return [int(rels[r]) for r in re.findall(r'<p:sldId id="\d+" r:id="([^"]+)"/>', lst)]


def _shapes(z, sn):
    """장 sn 의 화면 도형 텍스트 전부, 등장 순서대로.

    🟥 **첫 항목을 자르지 않는다.** 목차형 장(제목 도형 하나뿐)은 첫 도형을 빼면 화면이
    0개가 된다 — 신호 §3-4 실사고. `SKIP_SHAPES` 로 걸러진 것을 뺀 **나머지 전부**가 화면이다."""
    x = z.read('ppt/slides/slide%d.xml' % sn).decode('utf-8')
    out = []
    for m in re.finditer(r'<p:(sp|graphicFrame)>.*?</p:\1>', x, re.S):
        b = m.group(0)
        nm = re.search(r'name="([^"]*)"', b)
        nm = nm.group(1) if nm else '?'
        if nm in SKIP_SHAPES:
            continue
        t = ' '.join(re.findall(r'<a:t>([^<]*)</a:t>', b)).strip()
        if t:
            out.append((nm, re.sub(r'\s+', ' ', t)))
    return out


def _notes(z, sn):
    p = 'ppt/slides/_rels/slide%d.xml.rels' % sn
    if p not in z.namelist():
        return ''
    m = re.search(r'notesSlide(\d+)', z.read(p).decode('utf-8'))
    if not m:
        return ''
    npath = 'ppt/notesSlides/notesSlide%s.xml' % m.group(1)
    if npath not in z.namelist():
        return ''
    t = ' '.join(re.findall(r'<a:t>([^<]*)</a:t>', z.read(npath).decode('utf-8')))
    return re.sub(r'\s+', ' ', t).strip()


def build(deck_path):
    """계산부 — 디스크에 안 쓴다. (manuscript_lines, ids) 를 돌려준다.
    known-pair 셀프테스트가 파일 없이 이 함수를 직접 부른다."""
    z = zipfile.ZipFile(deck_path)
    order = _order(z)
    lines, ids = [], []
    for i, sn in enumerate(order, 1):
        sid = 'S%d' % i
        ids.append(sid)
        sh = _shapes(z, sn)
        head = (sh[0][1] if sh else '(제목 없음)')[:50]
        lines.append('### %s · %s' % (sid, head))
        lines.append('🖥')
        for _nm, t in sh:            # 🖥 전체 — 첫 도형 포함, 절대 자르지 않는다
            lines.append(t)
        lines.append('🗣')
        lines.append(_notes(z, sn) or '(발화 없음)')
        lines.append('')
    return lines, ids


def write(deck_path, out_dir):
    os.makedirs(out_dir, exist_ok=True)
    lines, ids = build(deck_path)
    open(os.path.join(out_dir, 'manuscript.md'), 'w', encoding='utf-8').write('\n'.join(lines))

    with open(os.path.join(out_dir, 'note_map.py'), 'w', encoding='utf-8') as f:
        f.write('# 생성물 — preprep_wire.py 가 만든다. 손으로 고치지 마라(재실행이 덮어쓴다).\n')
        f.write('# 이 어댑터는 절과 장을 1:1 로 배선한다(pptx 자체가 정본인 자산군 전제).\n')
        f.write('NOTE_MAP = {\n')
        for k, sid in enumerate(ids, 1):
            f.write('    %d: (%r, None),\n' % (k, sid))
        f.write('}\n')

    # 🟥 canon_terms/jargon_terms 는 basename 만 적는다 — 2026-09-04 수리로 root(<out_dir>)
    #   기준을 먼저 보고, 없으면 스킬 위치(HERE)로 자동 폴백하므로 example 을 그대로 가리켜도
    #   죽지 않는다(그 자리 대신 「(skill-dir fallback)」 표기가 뜬다). 실전 목록은 사람이
    #   <out_dir>/ 아래에 같은 이름으로 두면 root 쪽이 우선한다.
    surf = '''\
asset_family: "pptx-only 자산군 — preprep_wire 배선"
root: "."
canon_ledger: "CANON_LEDGER.md"           # 🟥 사람이 채운다. 없으면 L1-a 는 UNCOVERED 로 뜬다
canon_terms:  "canon_terms.example.yaml"  # 🟥 root 기준 — 여기 없으면 스킬 위치로 폴백
jargon_terms: "jargon_terms.example.yaml"

interslide:
  manuscript: script
  builder: "note_map.py"

surfaces:
  - id: script
    path: "manuscript.md"
    kind: markdown
    canonical_for: [발화 문자열, 화면 문자열]
    unit_pattern: '^###\\\\s+(S[0-9]+)\\\\s*·'
    speech_marker: "🗣"
    spoken: true
    retire_markers: ["🚫", "제거됨", "폐기"]

adjacent_dup:
  max_dup_chars: 40
'''
    open(os.path.join(out_dir, 'surfaces.yaml'), 'w', encoding='utf-8').write(surf)
    print('세움: %s/manuscript.md (%d절) · note_map.py · surfaces.yaml' % (out_dir, len(ids)))
    print('🟥 canon_terms/jargon_terms 는 example 을 가리킨다 — 실전은 사람이 %s 에 자기'
          ' 코퍼스의 금지어/조어 목록을 채워라(파일명을 같게 두면 root 쪽이 우선한다).'
          % out_dir)


def selftest():
    """known-pair — 신호 §3-4 «첫 도형 유실» 재발 방지. python-pptx 로 최소 덱을 지어 돈다."""
    import tempfile
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError:
        print('preprep_wire SELFTEST: SKIPPED — python-pptx 없음 (0 아님, UNMEASURED)')
        return 2

    ok = True
    with tempfile.TemporaryDirectory() as d:
        deck = os.path.join(d, 't.pptx')
        prs = Presentation()
        blank = prs.slide_layouts[6]

        # known-positive: 제목 도형 «하나뿐»인 목차형 장 — 화면이 비면 재발이다
        s1 = prs.slides.add_slide(blank)
        tb1 = s1.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        tb1.text_frame.text = '목차'

        # known-negative: 도형 둘(제목+본문) — 둘 다 화면에 남아야 정상
        s2 = prs.slides.add_slide(blank)
        tb2a = s2.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        tb2a.text_frame.text = '본문 제목'
        tb2b = s2.shapes.add_textbox(Inches(1), Inches(2), Inches(4), Inches(1))
        tb2b.text_frame.text = '본문 내용'

        prs.save(deck)
        lines, _ids = build(deck)
        text = '\n'.join(lines)
        i1, i2 = text.index('### S1'), text.index('### S2')
        block1, block2 = text[i1:i2], text[i2:]

        screen1 = block1.split('🖥', 1)[1].split('🗣', 1)[0]
        has1 = '목차' in screen1
        print(f'  {"PASS" if has1 else "FAIL"}  known-positive(목차형 장, 제목 도형 하나뿐): '
              f'화면에 «목차» {"있음" if has1 else "없음 — 첫 도형 유실 재발!"}')
        ok &= has1

        screen2 = block2.split('🖥', 1)[1].split('🗣', 1)[0]
        has2 = ('본문 제목' in screen2) and ('본문 내용' in screen2)
        print(f'  {"PASS" if has2 else "FAIL"}  known-negative(다중 도형 장, 둘 다 남아야 함): '
              f'{"둘 다 있음" if has2 else "유실됨"}')
        ok &= has2

    print('preprep_wire SELFTEST:', 'PASS' if ok else 'FAIL — 첫 도형 유실이 재발했다')
    return 0 if ok else 1


def main():
    if '--self-test' in sys.argv:
        return selftest()
    if len(sys.argv) < 2:
        print(__doc__)
        return 2
    deck = sys.argv[1]
    out = sys.argv[sys.argv.index('--out') + 1] if '--out' in sys.argv else '_preprep'
    write(deck, out)
    return 0


if __name__ == '__main__':
    sys.exit(main())
